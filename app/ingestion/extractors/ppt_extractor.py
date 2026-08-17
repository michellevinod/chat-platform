from pathlib import Path

from pptx import Presentation

from app.ingestion.extractors.base_extractor import BaseExtractor
from app.ingestion.extractors.raw_models import (
    RawBlock,
    RawDocument,
    RawPage,
    RawTableBlock,
    RawTextBlock,
)
from app.models.enums import BlockType


class PPTXExtractor(BaseExtractor):
    """
    Extracts text and tables from Microsoft PowerPoint (.pptx) files.
    Each slide is treated as one page.
    """

    def extract(
        self,
        file_path: Path,
    ) -> RawDocument:

        presentation = Presentation(file_path)

        pages: list[RawPage] = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):

            blocks: list[RawBlock] = []

            block_number = 0

            for shape in slide.shapes:

                if shape.has_table:
                    table = shape.table
                    rows_data = []
                    for row in table.rows:
                        row_vals = [cell.text.strip() for cell in row.cells]
                        rows_data.append(row_vals)

                    if rows_data and len(rows_data) > 0:
                        headers = rows_data[0]
                        lines = []
                        lines.append("| " + " | ".join(headers) + " |")
                        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                        for r in rows_data[1:]:
                            lines.append("| " + " | ".join(r) + " |")
                        md_table = "\n".join(lines)

                        blocks.append(
                            RawTableBlock(
                                page_number=slide_number,
                                block_number=block_number,
                                block_type=BlockType.TABLE,
                                bbox=(0.0, 0.0, 0.0, 0.0),
                                markdown=md_table,
                                rows=rows_data,
                                headers=headers,
                                caption=f"Slide {slide_number} Table",
                            )
                        )
                        block_number += 1
                    continue

                if not hasattr(shape, "text"):
                    continue

                text = shape.text.strip()

                if not text:
                    continue

                blocks.append(
                    RawTextBlock(
                        page_number=slide_number,
                        block_number=block_number,
                        block_type=BlockType.TEXT,
                        bbox=(0.0, 0.0, 0.0, 0.0),
                        text=text,
                    )
                )

                block_number += 1

            pages.append(
                RawPage(
                    page_number=slide_number,
                    blocks=blocks,
                )
            )

        return RawDocument(
            file_name=file_path.name,
            total_pages=len(pages),
            pages=pages,
        )